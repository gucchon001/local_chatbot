# document_processor.py
import os
import pandas as pd
from datetime import datetime
import docx
import xlrd
import csv
from pptx import Presentation
from langchain_community.document_loaders import PyPDFLoader, TextLoader
from langchain.text_splitter import CharacterTextSplitter
from langchain.schema import Document
import logging

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

def find_documents(folder_path):
    document_files = []
    for root, _, files in os.walk(folder_path):
        for file in files:
            if file.lower().endswith(('.pdf', '.xlsx', '.xls', '.docx', '.doc', '.txt', 'pptx', '.csv')):
                document_files.append(os.path.join(root, file))
    return document_files

def check_existing_data(config):
    parquet_file = config['parquet_file']
    faiss_index_file = config['faiss_index_file']
    
    if os.path.exists(parquet_file) and os.path.exists(faiss_index_file):
        parquet_mtime = datetime.fromtimestamp(os.path.getmtime(parquet_file))
        return True, parquet_mtime.strftime("%Y-%m-%d %H:%M:%S")
    return False, None

def analyze_documents(folder_path):
    logger.info(f"文書分析を開始: {folder_path}")
    document_files = find_documents(folder_path)
    
    stats = {
        "最終更新日": "情報なし",
        "ファイル数": 0,
        "総サイズ": 0,
        "ファイルタイプ": {}
    }

    if not document_files:
        logger.warning(f"ドキュメントが見つかりませんでした: {folder_path}")
        stats["警告"] = "ドキュメントが見つかりませんでした"
        return stats

    latest_update = datetime.min

    for file in document_files:
        file_extension = os.path.splitext(file)[1].lower()
        stats["ファイルタイプ"][file_extension] = stats["ファイルタイプ"].get(file_extension, 0) + 1
        file_size = os.path.getsize(file)
        stats["総サイズ"] += file_size
        file_update_time = datetime.fromtimestamp(os.path.getmtime(file))
        if file_update_time > latest_update:
            latest_update = file_update_time
        
        logger.debug(f"ファイル: {file}, サイズ: {file_size}, 更新日時: {file_update_time}")

    stats["最終更新日"] = latest_update.strftime("%Y-%m-%d %H:%M:%S")
    stats["ファイル数"] = len(document_files)
    
    logger.info(f"文書分析完了。統計情報: {stats}")
    return stats

def process_csv(file_path):
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            reader = csv.reader(file)
            text = "\n".join([",".join(row) for row in reader])
        return [Document(page_content=text, metadata={"source": file_path})]
    except UnicodeDecodeError:
        try:
            with open(file_path, 'r', encoding='shift_jis') as file:
                reader = csv.reader(file)
                text = "\n".join([",".join(row) for row in reader])
            return [Document(page_content=text, metadata={"source": file_path})]
        except Exception as e:
            logger.error(f"CSVファイルの読み込みに失敗しました: {file_path}, エラー: {e}")
            return None

def process_text(file_path):
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            text = file.read()
    except UnicodeDecodeError:
        try:
            with open(file_path, 'r', encoding='shift_jis') as file:
                text = file.read()
        except Exception as e:
            logger.error(f"テキストファイルの読み込みに失敗しました: {file_path}, エラー: {e}")
            return None
    return [Document(page_content=text, metadata={"source": file_path})]

def process_excel(file_path):
    try:
        df = pd.read_excel(file_path, engine='openpyxl')
        text = df.to_string(index=False)
    except Exception as e:
        logger.warning(f"openpyxlでの読み込みに失敗しました: {file_path}, エラー: {e}")
        try:
            workbook = xlrd.open_workbook(file_path)
            text = ""
            for sheet in workbook.sheets():
                for row in range(sheet.nrows):
                    text += " ".join(str(sheet.cell_value(row, col)) for col in range(sheet.ncols)) + "\n"
        except Exception as e:
            logger.error(f"xlrdでの読み込みにも失敗しました: {file_path}, エラー: {e}")
            return None
    return [Document(page_content=text, metadata={"source": file_path})]

def process_word(file_path):
    try:
        doc = docx.Document(file_path)
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        return [Document(page_content=text, metadata={"source": file_path})]
    except Exception as e:
        logger.error(f"Wordファイルの読み込みに失敗しました: {file_path}, エラー: {e}")
        return None

def process_pptx(file_path):
    try:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return [Document(page_content=text, metadata={"source": file_path})]
    except Exception as e:
        logger.error(f"PowerPointファイルの読み込みに失敗しました: {file_path}, エラー: {e}")
        return None

def process_document(file_path):
    file_extension = os.path.splitext(file_path)[1].lower()
    
    try:
        if file_extension in ['.pdf']:
            loader = PyPDFLoader(file_path)
            documents = loader.load()
        elif file_extension in ['.xlsx', '.xls']:
            documents = process_excel(file_path)
        elif file_extension in ['.docx', '.doc']:
            documents = process_word(file_path)
        elif file_extension in ['.txt']:
            documents = process_text(file_path)
        elif file_extension in ['.pptx']:
            documents = process_pptx(file_path)
        elif file_extension in ['.csv']:
            documents = process_csv(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_extension}")

        if documents is None:
            return []

        text_splitter = CharacterTextSplitter(chunk_size=1000, chunk_overlap=200, separator="\n")
        return text_splitter.split_documents(documents)
    except Exception as e:
        logger.error(f"ファイル {file_path} の処理中にエラーが発生しました: {e}")
        return []

def process_all_documents(folder_path):
    logger.info(f"文書処理を開始: {folder_path}")
    stats, document_files = analyze_documents(folder_path)
    
    logger.info(f"処理対象ファイル数: {len(document_files)}")
    logger.info(f"推定処理時間: {len(document_files) * 2} 秒") # 1ファイルあたり2秒と仮定
    
    all_chunks = []
    for doc_file in document_files:
        logger.info(f"処理中: {doc_file}")
        chunks = process_document(doc_file)
        all_chunks.extend(chunks)
        logger.info(f"処理完了: {doc_file}, チャンク数: {len(chunks)}")
    
    logger.info(f"全文書処理完了。総チャンク数: {len(all_chunks)}")
    return all_chunks, stats

def get_file_statistics(folder_path):
    logger.info(f"get_file_statistics が呼び出されました: {folder_path}")
    return analyze_documents(folder_path)

def process_changed_documents(folder_path, changed_files):
    logger.info(f"変更された文書の処理を開始: {folder_path}")
    logger.info(f"変更されたファイル数: {len(changed_files)}")
    
    changed_chunks = []
    for doc_file in changed_files:
        logger.info(f"変更ファイルを処理中: {doc_file}")
        chunks = process_document(doc_file)
        changed_chunks.extend(chunks)
        logger.info(f"変更ファイル処理完了: {doc_file}, チャンク数: {len(chunks)}")
    
    logger.info(f"変更された全文書の処理完了。総チャンク数: {len(changed_chunks)}")
    return changed_chunks